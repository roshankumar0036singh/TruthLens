from pptx import Presentation
from pptx.util import Inches, Pt, Emu

prs = Presentation(r'd:\DevYatra\TruthLens\DevYatra-Hackfest-template.pptx')

print(f'Slide Width: {prs.slide_width}, Height: {prs.slide_height}')
print(f'Slide Width (inches): {prs.slide_width / 914400:.2f}, Height (inches): {prs.slide_height / 914400:.2f}')
print(f'Total Slides: {len(prs.slides)}')
print()

for i, slide in enumerate(prs.slides):
    layout = slide.slide_layout
    print(f'=== SLIDE {i+1} ===')
    print(f'  Layout Name: {layout.name}')
    print(f'  Shapes count: {len(slide.shapes)}')
    for shape in slide.shapes:
        print(f'  Shape: {shape.shape_type}, Name: "{shape.name}"')
        print(f'    Position: left={shape.left}, top={shape.top}, width={shape.width}, height={shape.height}')
        if hasattr(shape, 'text') and shape.text:
            text_preview = shape.text[:200].replace('\n', ' | ')
            print(f'    Text: "{text_preview}"')
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for run in para.runs:
                    font = run.font
                    try:
                        color_str = str(font.color.rgb) if font.color and font.color.type else "inherit"
                    except:
                        color_str = "unknown"
                    print(f'    Run[{pi}]: size={font.size}, bold={font.bold}, color={color_str}, text="{run.text[:80]}"')
    print()

# Also check slide layouts available
print("=== AVAILABLE LAYOUTS ===")
for i, layout in enumerate(prs.slide_layouts):
    print(f"  Layout {i}: {layout.name}")
